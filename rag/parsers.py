"""Format-specific parsers that produce a normalized stream of Blocks.

Each parser returns List[Block]. A Block is the smallest unit of text the
chunker will consider for grouping; it carries provenance metadata so we can
build citations later (page number, slide number, heading path, sheet name).
"""
from __future__ import annotations

import re
from collections import Counter
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any


@dataclass
class Block:
    text: str
    kind: str  # "paragraph" | "heading" | "table_row" | "slide" | "sheet_row"
    meta: dict[str, Any] = field(default_factory=dict)


def parse(path: Path) -> list[Block]:
    ext = path.suffix.lower()
    if ext == ".pdf":
        return parse_pdf(path)
    if ext == ".docx":
        return parse_docx(path)
    if ext == ".pptx":
        return parse_pptx(path)
    if ext == ".xlsx":
        return parse_xlsx(path)
    if ext in (".txt", ".md"):
        return parse_text(path)
    raise ValueError(f"Unsupported file type: {ext}")


def parse_text(path: Path) -> list[Block]:
    text = path.read_text(encoding="utf-8", errors="replace")
    blocks: list[Block] = []
    for para in text.split("\n\n"):
        para = para.strip()
        if para:
            blocks.append(Block(text=para, kind="paragraph", meta={}))
    return blocks


def parse_pdf(path: Path) -> list[Block]:
    """Section-aware PDF parser.

    Walks each page's structured text (PyMuPDF's "dict" output) so we can read
    font sizes per block. Blocks whose max font size is meaningfully larger
    than the document's body font are treated as headings; we maintain a
    heading_stack so every emitted paragraph carries the section context it
    belongs to (same shape as parse_docx). The chunker then keeps each chunk
    within one section via heading_path boundary tracking.

    Falls back to legacy paragraph-mode parsing if the structural extraction
    can't determine a body font (very small / image-only / odd PDFs).
    """
    import fitz  # PyMuPDF

    from . import config_facade as cfg

    max_font = float(cfg.PDF_MAX_FONT_SIZE or 0)  # 0 = disabled
    screen_fonts = [
        f.strip().lower() for f in (cfg.PDF_SCREEN_FONTS or "").split(",") if f.strip()
    ]

    def _block_is_screen(sizes_spans_text: list[tuple[float, str]]) -> bool:
        """True if every span in the block matches a configured screen font."""
        if not screen_fonts or not sizes_spans_text:
            return False
        for _, font_name in sizes_spans_text:
            fn = (font_name or "").lower()
            if not any(sf in fn for sf in screen_fonts):
                return False
        return True

    doc = fitz.open(str(path))
    try:
        # First pass — plain text per page, used to identify repeating
        # headers/footers we should drop from the structured pass.
        pages_text: list[str] = []
        for page in doc:
            try:
                pages_text.append(page.get_text("text") or "")
            except Exception:
                pages_text.append("")

        top_repeats, bot_repeats = _find_repeating_lines(pages_text, min_frac=0.5)
        repeating = top_repeats | bot_repeats

        def is_repeating(text: str) -> bool:
            if not repeating:
                return False
            return re.sub(r"\d+", "#", text.strip()) in repeating

        # Estimate the document's body font size. If we can't, fall back.
        body_size = _estimate_body_font_size(doc)
        if body_size is None:
            return _parse_pdf_legacy(pages_text)

        heading_threshold = body_size * 1.15  # >= 15% larger than body = heading
        # (level, text) tuples so a new level-N heading correctly removes every
        # entry with level >= N — necessary for multiple level-5 field labels
        # to become siblings instead of nesting under each other, and for a
        # second Field Descriptions block to reset both slot 4 and slot 5.
        heading_stack: list[tuple[int, str]] = []

        def _stack_path() -> str:
            return " > ".join(t for _, t in heading_stack) if heading_stack else ""
        blocks: list[Block] = []
        # Buffer for consecutive same-page screen blocks so they merge into one.
        screen_buf: list[str] = []
        screen_buf_page: int | None = None
        # Tracks the most recent screen so subsequent blocks can point back to
        # it. Reset by major section boundaries (H1/H2/H3).
        last_screen_page: int | None = None

        def _flush_screen() -> None:
            nonlocal screen_buf, screen_buf_page, last_screen_page
            if not screen_buf or screen_buf_page is None:
                return
            joined = "\n".join(screen_buf)
            blocks.append(Block(
                text=joined,
                kind="screen",
                meta={
                    "page": screen_buf_page,
                    "heading_path": _stack_path(),
                },
            ))
            last_screen_page = screen_buf_page
            screen_buf = []
            screen_buf_page = None

        for page_num, page in enumerate(doc, start=1):
            if screen_buf_page is not None and screen_buf_page != page_num:
                _flush_screen()
            try:
                page_dict = page.get_text("dict")
            except Exception:
                continue
            for blk in page_dict.get("blocks", []):
                if blk.get("type", 0) != 0:  # 0 = text block
                    continue

                # Process lines individually so a block containing a heading
                # line followed by body lines (common on page transitions where
                # PyMuPDF glues them together) still gets split cleanly.
                para_lines: list[str] = []
                para_max_size: float = 0.0

                def _flush_paragraph() -> None:
                    nonlocal para_lines, para_max_size, screen_buf
                    if not para_lines:
                        return
                    ptext = " ".join(para_lines).strip()
                    para_lines = []
                    p_max = para_max_size
                    para_max_size = 0.0
                    if not ptext or is_repeating(ptext):
                        return
                    if max_font and p_max >= max_font:
                        return
                    if screen_buf:
                        _flush_screen()
                    meta = {
                        "page": page_num,
                        "heading_path": _stack_path(),
                    }
                    if last_screen_page is not None:
                        meta["related_screen_page"] = last_screen_page
                    blocks.append(Block(text=ptext, kind="paragraph", meta=meta))

                for line in blk.get("lines", []):
                    spans = line.get("spans", [])
                    line_text = "".join(s.get("text", "") for s in spans).strip()
                    if not line_text or is_repeating(line_text):
                        continue

                    line_sizes: list[float] = []
                    line_span_fonts: list[tuple[float, str]] = []
                    line_bold = 0
                    line_total = 0
                    for s in spans:
                        sz = s.get("size")
                        if sz:
                            line_sizes.append(sz)
                        line_span_fonts.append((sz or 0.0, s.get("font") or ""))
                        span_text = s.get("text") or ""
                        span_chars = len(span_text)
                        line_total += span_chars
                        if int(s.get("flags") or 0) & 16:
                            line_bold += span_chars
                    line_max_size = max(line_sizes) if line_sizes else body_size

                    # Screen detection at line level. Consecutive same-page
                    # screen lines still buffer into one Block(kind="screen").
                    if _block_is_screen(line_span_fonts):
                        _flush_paragraph()
                        if screen_buf_page is not None and screen_buf_page != page_num:
                            _flush_screen()
                        screen_buf_page = page_num
                        screen_buf.append(line_text)
                        continue
                    # Non-screen line — flush any pending screen first.
                    if screen_buf:
                        _flush_screen()

                    if max_font and line_max_size >= max_font:
                        _flush_paragraph()
                        continue

                    line_is_bold = line_total > 0 and (line_bold / line_total) >= 0.5
                    is_heading_by_size = line_max_size >= heading_threshold and len(line_text) <= 200
                    is_heading_by_bold = line_is_bold and line_max_size >= body_size and len(line_text) <= 100
                    is_heading = is_heading_by_size or is_heading_by_bold

                    if is_heading:
                        # Flush any accumulated body lines as their own
                        # paragraph before emitting the heading.
                        _flush_paragraph()
                        level = _heading_level_from_size(line_max_size, body_size, line_is_bold)
                        if level <= 3:
                            last_screen_page = None
                        # Drop every stack entry at this level or deeper; a
                        # level-N heading is a peer to prior level-N entries,
                        # not their child.
                        heading_stack = [(l, t) for l, t in heading_stack if l < level]
                        heading_stack.append((level, line_text))
                        meta = {
                            "page": page_num,
                            "heading_path": _stack_path(),
                            "level": level,
                        }
                        if last_screen_page is not None:
                            meta["related_screen_page"] = last_screen_page
                        blocks.append(Block(text=line_text, kind="heading", meta=meta))
                    else:
                        para_lines.append(line_text)
                        if line_max_size > para_max_size:
                            para_max_size = line_max_size

                # End of block — flush any accumulated paragraph.
                _flush_paragraph()
        _flush_screen()  # emit any trailing screen block
        return blocks
    finally:
        doc.close()


def _parse_pdf_legacy(pages_text: list[str]) -> list[Block]:
    """Paragraph-mode fallback (the pre-section-aware behavior)."""
    pages_text = _strip_repeating_lines(pages_text, min_frac=0.5)
    blocks: list[Block] = []
    for page_num, text in enumerate(pages_text, start=1):
        for para in _split_paragraphs(text):
            blocks.append(Block(text=para, kind="paragraph", meta={"page": page_num}))
    return blocks


def _estimate_body_font_size(doc) -> float | None:
    """Most-common font size weighted by character count. Returns None if
    no usable text/font info is available."""
    sizes: Counter[float] = Counter()
    for page in doc:
        try:
            page_dict = page.get_text("dict")
        except Exception:
            continue
        for blk in page_dict.get("blocks", []):
            if blk.get("type", 0) != 0:
                continue
            for line in blk.get("lines", []):
                for span in line.get("spans", []):
                    sz = span.get("size")
                    txt = (span.get("text") or "").strip()
                    if sz and txt:
                        # Round to nearest 0.5pt so near-identical sizes group.
                        sizes[round(sz * 2) / 2] += len(txt)
    if not sizes:
        return None
    return sizes.most_common(1)[0][0]


def _heading_level_from_size(size: float, body_size: float, is_bold: bool = False) -> int:
    """Map font-size ratio (plus bold flag) to a heading level (1 = biggest).

    Levels 1-3 are size-based. Levels 4-5 are the bold-only regime for headings
    that are at or near body size (e.g., 'Field Descriptions' at 11pt on 10pt
    body, or field-name labels at 10pt bold on 10pt body). Keeping them at
    distinct levels lets each field name nest under its parent 'Field
    Descriptions' block instead of appearing as its sibling.
    """
    ratio = size / body_size if body_size else 1.0
    if ratio >= 2.0:
        return 1
    if ratio >= 1.5:
        return 2
    if ratio >= 1.15:
        return 3
    # Bold-only regime.
    if is_bold and ratio >= 1.05:
        return 4
    return 5


def parse_docx(path: Path) -> list[Block]:
    from docx import Document

    from . import config_facade as cfg

    doc = Document(str(path))
    blocks: list[Block] = []
    heading_stack: list[str] = []  # [h1, h2, h3, ...]
    stop_keywords = [k.lower() for k in cfg.DOCX_STOP_HEADINGS if k.strip()]

    para_map = {p._element: p for p in doc.paragraphs}
    table_map = {t._element: t for t in doc.tables}

    for child in doc.element.body.iterchildren():
        tag = child.tag.split("}", 1)[-1]

        if tag == "p":
            para = para_map.get(child)
            if para is None:
                continue
            text = (para.text or "").strip()
            if not text:
                continue
            style = (para.style.name or "").lower() if para.style else ""
            is_heading = style.startswith("heading")

            if is_heading and stop_keywords:
                heading_lower = text.lower()
                if any(kw in heading_lower for kw in stop_keywords):
                    break  # drop the matching heading and everything after

            if is_heading:
                try:
                    level = int(style.split()[-1])
                except ValueError:
                    level = 1
                heading_stack = heading_stack[: max(0, level - 1)]
                heading_stack.append(text)
                blocks.append(Block(
                    text=text,
                    kind="heading",
                    meta={"heading_path": " > ".join(heading_stack), "level": level},
                ))
            else:
                blocks.append(Block(
                    text=text,
                    kind="paragraph",
                    meta={"heading_path": " > ".join(heading_stack) if heading_stack else ""},
                ))

        elif tag == "tbl":
            table = table_map.get(child)
            if table is None:
                continue
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells]
                row_text = " | ".join(c for c in cells if c)
                if row_text:
                    blocks.append(Block(
                        text=row_text,
                        kind="table_row",
                        meta={"heading_path": " > ".join(heading_stack) if heading_stack else ""},
                    ))

    return blocks


def parse_pptx(path: Path) -> list[Block]:
    from pptx import Presentation

    prs = Presentation(str(path))
    blocks: list[Block] = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        title = ""
        body_parts: list[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            if shape == slide.shapes.title or (
                shape.has_text_frame and shape.text_frame.text and not title
            ):
                if not title:
                    title = text
                    continue
            body_parts.append(text)

        notes = ""
        if slide.has_notes_slide:
            notes = (slide.notes_slide.notes_text_frame.text or "").strip()

        combined = "\n\n".join(part for part in [title, *body_parts, notes] if part)
        if combined:
            blocks.append(Block(
                text=combined,
                kind="slide",
                meta={"slide_num": slide_num, "slide_title": title},
            ))
    return blocks


def parse_xlsx(path: Path) -> list[Block]:
    from openpyxl import load_workbook

    wb = load_workbook(str(path), data_only=True, read_only=True)
    blocks: list[Block] = []
    for sheet in wb.worksheets:
        rows_iter = sheet.iter_rows(values_only=True)
        try:
            header = next(rows_iter)
        except StopIteration:
            continue
        header_str = " | ".join(str(c) if c is not None else "" for c in header)
        buffer: list[str] = []
        first_row_in_buffer = 2
        row_idx = 2
        for row in rows_iter:
            cells = " | ".join(str(c) if c is not None else "" for c in row)
            if cells.strip():
                buffer.append(cells)
            if len(buffer) >= 20:
                blocks.append(Block(
                    text=f"Headers: {header_str}\n" + "\n".join(buffer),
                    kind="sheet_row",
                    meta={
                        "sheet": sheet.title,
                        "row_start": first_row_in_buffer,
                        "row_end": row_idx,
                    },
                ))
                buffer = []
                first_row_in_buffer = row_idx + 1
            row_idx += 1
        if buffer:
            blocks.append(Block(
                text=f"Headers: {header_str}\n" + "\n".join(buffer),
                kind="sheet_row",
                meta={
                    "sheet": sheet.title,
                    "row_start": first_row_in_buffer,
                    "row_end": row_idx - 1,
                },
            ))
    wb.close()
    return blocks


def _split_paragraphs(text: str) -> list[str]:
    parts = [p.strip() for p in text.split("\n\n")]
    parts = [p for p in parts if p]
    if not parts:
        parts = [line.strip() for line in text.splitlines() if line.strip()]
    return parts


def _find_repeating_lines(
    pages_text: list[str], min_frac: float = 0.5
) -> tuple[set[str], set[str]]:
    """Identify digit-normalized header/footer lines that recur on most pages.

    Returns (top_repeats, bot_repeats) — sets of digit-normalized strings.
    Needs at least 3 pages to engage; returns (set(), set()) otherwise.
    """
    if len(pages_text) < 3:
        return set(), set()

    def norm(line: str) -> str:
        return re.sub(r"\d+", "#", line.strip())

    page_lines = [p.splitlines() for p in pages_text]
    top_counts: Counter[str] = Counter()
    bot_counts: Counter[str] = Counter()
    for lines in page_lines:
        non_empty = [l for l in lines if l.strip()]
        for l in non_empty[:2]:
            top_counts[norm(l)] += 1
        for l in non_empty[-2:]:
            bot_counts[norm(l)] += 1

    threshold = max(2, int(min_frac * len(pages_text)))
    top_repeats = {k for k, v in top_counts.items() if v >= threshold}
    bot_repeats = {k for k, v in bot_counts.items() if v >= threshold}
    return top_repeats, bot_repeats


def _strip_repeating_lines(pages_text: list[str], min_frac: float = 0.5) -> list[str]:
    """Remove running headers/footers that repeat at the top or bottom of pages."""
    top_repeats, bot_repeats = _find_repeating_lines(pages_text, min_frac)
    if not top_repeats and not bot_repeats:
        return pages_text

    def norm(line: str) -> str:
        return re.sub(r"\d+", "#", line.strip())

    page_lines = [p.splitlines() for p in pages_text]
    cleaned: list[str] = []
    for lines in page_lines:
        i, j = 0, len(lines)
        while i < j and (not lines[i].strip() or norm(lines[i]) in top_repeats):
            i += 1
        while j > i and (not lines[j - 1].strip() or norm(lines[j - 1]) in bot_repeats):
            j -= 1
        cleaned.append("\n".join(lines[i:j]))
    return cleaned
